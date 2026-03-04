#!/usr/bin/env python3
"""
生成中期答辩PPT v2 - 修复排版
- 用"空白"布局消除占位符残留
- 内容铺满全页，合理间距
- Mermaid渲染PNG嵌入架构/流程图
- Kalman合并到控制算法页，总计19页
"""

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import os

# ============================================================
# 配置
# ============================================================
VERSION = "v3"
TEMPLATE_PATH = "yyh_中期答辩.pptx"
OUTPUT_PATH = f"yyh_中期答辩_{VERSION}.pptx"
MERMAID_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_images")

# 主题色
C_PRI = RGBColor(0x2B, 0x75, 0x85)
C_DARK = RGBColor(0x1A, 0x4A, 0x55)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x33, 0x33, 0x33)
C_GRAY = RGBColor(0x66, 0x66, 0x66)
C_GREEN = RGBColor(0x27, 0xAE, 0x60)
C_YELLOW = RGBColor(0xF3, 0x9C, 0x12)
C_RED = RGBColor(0xE7, 0x4C, 0x3C)
C_LIGHT = RGBColor(0xF0, 0xF8, 0xFA)

# 字体
F_TITLE = "微软雅黑"
F_BODY = "微软雅黑"
F_CODE = "Consolas"
F_MATH = "Times New Roman"

# 页面尺寸
SW = 12192000  # slide width EMU
SH = 6858000   # slide height EMU
ML = Cm(1.2)   # margin left
MR = Cm(1.2)   # margin right
CW = SW - ML - MR  # content width

# 背景图位置 (从原始封面)
BG_L, BG_T, BG_W, BG_H = 995680, 2978785, 9846310, 4135755


# ============================================================
# 辅助函数
# ============================================================
def sf(run, name=F_BODY, size=Pt(14), bold=False, color=C_BLACK, italic=False):
    """set font"""
    run.font.name = name
    run.font.east_asian = "微软雅黑"
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic


def tb(slide, l, t, w, h, text="", fn=F_BODY, fs=Pt(14), bold=False,
       color=C_BLACK, align=PP_ALIGN.LEFT):
    """add textbox"""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    sf(r, fn, fs, bold, color)
    return box


def title(slide, text, y=Cm(0.8)):
    """页面标题 - 大字加粗"""
    return tb(slide, ML, y, CW, Cm(1.2), text, F_TITLE, Pt(28), True, C_PRI)


def bullets(slide, l, t, w, h, items, fs=Pt(13), color=C_BLACK, sp=Pt(5), bc="•"):
    """带要点的文本框"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = sp
        r = p.add_run()
        r.text = f"{bc} {item}"
        sf(r, F_BODY, fs, False, color)
    return box


def tbl(slide, l, t, w, h, data, cw=None, fs=Pt(12), hc=C_PRI):
    """添加表格"""
    nr, nc = len(data), len(data[0])
    sh = slide.shapes.add_table(nr, nc, l, t, w, h)
    table = sh.table
    if cw:
        for i, ww in enumerate(cw):
            table.columns[i].width = ww
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = F_BODY
                    run.font.east_asian = F_BODY
                    run.font.size = fs
                    run.font.bold = r == 0
                    run.font.color.rgb = C_WHITE if r == 0 else C_BLACK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hc
    return sh


def rrect(slide, l, t, w, h, fill=C_LIGHT, text="", fs=Pt(12), fc=C_BLACK,
          bold=False, border=C_PRI, bw=Pt(1)):
    """圆角矩形"""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = bw
    else:
        sh.line.fill.background()
    if text:
        sh.text_frame.word_wrap = True
        sh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = sh.text_frame.paragraphs[0].add_run()
        r.text = text
        sf(r, F_BODY, fs, bold, fc)
    return sh


def placeholder(slide, l, t, w, h, label="[图片占位]"):
    """图片占位灰框"""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    sh.line.color.rgb = C_GRAY
    sh.line.width = Pt(1)
    sh.line.dash_style = 2
    sh.text_frame.word_wrap = True
    p = sh.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    sf(r, F_BODY, Pt(12), False, C_GRAY)
    return sh


def mpara(slide, l, t, w, h, paras):
    """多段落文本框. paras: list of dict{text, fs, bold, color, font, align}"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, pd in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pd.get("align", PP_ALIGN.LEFT)
        p.space_after = pd.get("sa", Pt(3))
        r = p.add_run()
        r.text = pd["text"]
        sf(r, pd.get("font", F_BODY), pd.get("fs", Pt(13)),
           pd.get("bold", False), pd.get("color", C_BLACK))
    return box


def add_bg(slide):
    """添加背景底纹(校园线框图)"""
    bg_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_images", "bg.png")
    if not os.path.exists(bg_path):
        bg_path = "/tmp/ppt_bg.png"  # fallback
    if os.path.exists(bg_path):
        pic = slide.shapes.add_picture(bg_path, BG_L, BG_T, BG_W, BG_H)
        sp = pic._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)


def add_mermaid(slide, name, l, t, w, h):
    """嵌入mermaid渲染的图片，保持宽高比居中放置，失败则用占位框"""
    path = os.path.join(MERMAID_DIR, f"{name}.jpg")
    if not os.path.exists(path):
        placeholder(slide, l, t, w, h, f"[Mermaid: {name}]")
        return
    from PIL import Image
    img = Image.open(path)
    iw, ih = img.size  # pixels
    # 计算适配尺寸（保持宽高比，fit within w×h）
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        # 图片更宽 → 以宽度为准
        final_w = w
        final_h = int(w / img_ratio)
    else:
        # 图片更高 → 以高度为准
        final_h = h
        final_w = int(h * img_ratio)
    # 居中偏移
    offset_l = l + (w - final_w) // 2
    offset_t = t + (h - final_h) // 2
    slide.shapes.add_picture(path, offset_l, offset_t, final_w, final_h)


def section_tag(slide, text):
    """右上角节标签"""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 SW - Cm(6.5), Cm(0.2), Cm(6), Cm(0.55))
    sh.fill.solid()
    sh.fill.fore_color.rgb = C_PRI
    sh.line.fill.background()
    r = sh.text_frame.paragraphs[0].add_run()
    r.text = text
    sh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    sf(r, F_BODY, Pt(10), True, C_WHITE)


def page_num(slide, n, total=19):
    """页码"""
    tb(slide, SW - Cm(2), SH - Cm(0.8), Cm(1.5), Cm(0.5),
       f"{n}/{total}", F_BODY, Pt(9), False, C_GRAY, PP_ALIGN.RIGHT)


def blank(prs):
    """创建空白幻灯片"""
    lm = {l.name: l for l in prs.slide_layouts}
    return prs.slides.add_slide(lm["空白"])


# ============================================================
# 主逻辑
# ============================================================
prs = Presentation(TEMPLATE_PATH)

# 删除第2页空白slide
if len(prs.slides) > 1:
    rId = prs.slides._sldIdLst[1].get(qn('r:id'))
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[1])

# --- 更新封面 (P1) ---
s1 = prs.slides[0]
for shape in s1.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if "Presenters" in para.text:
                for run in para.runs:
                    run.text = "答辩人: 杨倚航    学号: 12211113    导师: 柯文德"
            if "2026/3/4" in para.text:
                for run in para.runs:
                    run.text = "南方科技大学 · 2026年3月"


# ==========================================================
# P2 - 研究背景
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "研究背景")
section_tag(s, "Part 1 · 背景与目标")
page_num(s, 2)

bullets(s, ML, Cm(2.5), Cm(16), Cm(8), [
    "赛题需求: RoboMaster工程机器人赛项需在3分钟内完成矿石采集→搬运→兑换全流程",
    "",
    "现有方案痛点:",
    "    ✗ 纯手动遥控6DOF机械臂，精度低",
    "    ✗ 操作手反应延迟，效率受限",
    "    ✗ 无法多任务并行",
    "",
    "本项目定位:",
    "    → 自主抓取物流机器人",
    "    → 视觉伺服闭环控制",
    "    → 感知→规划→执行 全链路自动化",
], Pt(14), C_BLACK, Pt(2))

placeholder(s, Cm(19), Cm(2.5), Cm(10.5), Cm(4), "[比赛场地照片]")
placeholder(s, Cm(19), Cm(7), Cm(10.5), Cm(3.5), "[矿石搬运示意图]")


# ==========================================================
# P3 - 研究目标与任务分解
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "研究目标与任务分解")
section_tag(s, "Part 1 · 背景与目标")
page_num(s, 3)

# 总目标
rrect(s, ML, Cm(2.3), CW, Cm(1.3), C_PRI,
      "总目标: 设计并实现一套基于ROS2的六轴机械臂视觉伺服控制系统，用于自主矿石搬运",
      Pt(14), C_WHITE, True, border=None)

# 四卡片
cw_ = Cm(7)
ch_ = Cm(5)
cy = Cm(4.2)
cards = [
    ("子任务1 ✅", "运动学建模", ["6DOF逆运动学", "D-H位姿解耦", "计算<50μs"], C_GREEN),
    ("子任务2 ✅", "动力学力矩控制", ["KDL动力学前馈", "级联PID控制", "速度后处理·200Hz"], C_GREEN),
    ("子任务3 ✅", "嵌入式系统开发", ["STM32+CAN总线", "FSM状态机·7状态", "Seasky协议·1000Hz"], C_GREEN),
    ("子任务4 🔧", "视觉伺服集成", ["YOLOv8目标检测", "PnP位姿估计", "闭环抓取(进行中)"], C_YELLOW),
]
for i, (tag, sub, items, color) in enumerate(cards):
    x = ML + i * (cw_ + Cm(0.3))
    r = rrect(s, x, cy, cw_, ch_, C_LIGHT, border=color, bw=Pt(2))
    tb(s, x + Cm(0.4), cy + Cm(0.2), cw_ - Cm(0.8), Cm(0.6),
       tag, F_BODY, Pt(12), True, color)
    tb(s, x + Cm(0.4), cy + Cm(0.9), cw_ - Cm(0.8), Cm(0.6),
       sub, F_BODY, Pt(14), True, C_DARK)
    bullets(s, x + Cm(0.4), cy + Cm(1.7), cw_ - Cm(0.8), Cm(2.5),
            items, Pt(11), C_GRAY, Pt(2), "·")

# 进度条
by = cy + ch_ + Cm(0.5)
bg_bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ML, by, CW, Cm(0.5))
bg_bar.fill.solid()
bg_bar.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
bg_bar.line.fill.background()
fg_bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ML, by, int(CW * 0.75), Cm(0.5))
fg_bar.fill.solid()
fg_bar.fill.fore_color.rgb = C_GREEN
fg_bar.line.fill.background()
tb(s, ML + int(CW * 0.75) + Cm(0.3), by - Cm(0.05),
   Cm(3), Cm(0.6), "进度 75%", F_BODY, Pt(11), True, C_GREEN)


# ==========================================================
# P4 - 系统总体架构 ★  (mermaid图)
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "系统总体架构 ★")
section_tag(s, "Part 2 · 系统架构")
page_num(s, 4)

add_mermaid(s, "system_arch", Cm(2), Cm(2.2), Cm(16), Cm(9))

# 右侧关键标注
mpara(s, Cm(19), Cm(2.5), Cm(10), Cm(8), [
    {"text": "关键数据流", "fs": Pt(15), "bold": True, "color": C_PRI},
    {"text": ""},
    {"text": "↓ /effort_controller/commands", "fs": Pt(11), "color": C_RED, "font": F_CODE},
    {"text": "   6轴力矩(Nm) + 1夹爪力(N) = 7元素", "fs": Pt(11)},
    {"text": "   200Hz 下发", "fs": Pt(11), "color": C_GRAY},
    {"text": ""},
    {"text": "↑ /joint_states", "fs": Pt(11), "color": RGBColor(0x27, 0x64, 0xBA), "font": F_CODE},
    {"text": "   7关节(位置+速度+存活标志)", "fs": Pt(11)},
    {"text": "   200Hz 上报", "fs": Pt(11), "color": C_GRAY},
    {"text": ""},
    {"text": "通信协议", "fs": Pt(14), "bold": True, "color": C_PRI},
    {"text": "· USB-UART 921600 baud", "fs": Pt(12)},
    {"text": "· Seasky自定义协议", "fs": Pt(12)},
    {"text": "· CRC8(头)+CRC16(全帧) 双校验", "fs": Pt(12)},
])


# ==========================================================
# P5 - 硬件平台
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "硬件平台")
section_tag(s, "Part 2 · 系统架构")
page_num(s, 5)

placeholder(s, ML, Cm(2.5), Cm(13), Cm(8.5),
            "[机器人实物照片 📷1/📷2\n用箭头标注J1~J6和夹爪位置]")

tbl(s, Cm(16), Cm(2.5), Cm(13), Cm(5.5),
    [["关节", "电机", "力矩限制"],
     ["J1", "J8009", "40 Nm"],
     ["J2", "J8009", "40 Nm"],
     ["J3", "J8009", "40 Nm"],
     ["J4", "GM6020", "1.2 Nm"],
     ["J5", "J4310", "7 Nm"],
     ["J6", "M2006", "1 Nm"],
     ["夹爪", "棱柱副", "70 N"]],
    [Cm(3), Cm(5), Cm(5)])

bullets(s, Cm(16), Cm(8.5), Cm(13), Cm(2.5), [
    "夹爪: 电机驱动棱柱副, 0~40mm行程, 峰值70N",
    "底盘: M3508×4 麦克纳姆轮",
    "主控: STM32F4 · 通信: USB-UART 921600",
], Pt(12), C_GRAY, Pt(3), "·")


# ==========================================================
# P6 - 软件技术栈
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "软件技术栈")
section_tag(s, "Part 2 · 系统架构")
page_num(s, 6)

layers = [
    ("应用层", "Mission Executor · Trajectory Manager · Cartesian Controller",
     RGBColor(0x42, 0xA5, 0xF5)),
    ("规划层", "MoveIt2 + Pilz Industrial Motion Planner (LIN/PTP)",
     RGBColor(0x66, 0xBB, 0x6A)),
    ("控制层", "KDL动力学 · 级联PID · 速度后处理滤波",
     RGBColor(0xFF, 0xA7, 0x26)),
    ("接口层", "MuJoCo + ImGui/ImPlot · Seasky串口 + serial_driver",
     RGBColor(0xAB, 0x47, 0xBC)),
    ("基础设施", "ROS2 Jazzy · Eigen3 · GLFW · yaml-cpp · ncurses",
     RGBColor(0x78, 0x90, 0x9C)),
]
lh = Cm(1.6)
lg = Cm(0.3)
lw = Cm(18)
for i, (name, desc, color) in enumerate(layers):
    y = Cm(2.5) + i * (lh + lg)
    r = rrect(s, ML, y, lw, lh, color, border=None)
    tb(s, ML + Cm(0.5), y + Cm(0.15), Cm(3.5), Cm(0.6),
       name, F_BODY, Pt(13), True, C_WHITE)
    tb(s, ML + Cm(4.5), y + Cm(0.15), Cm(13), Cm(1.2),
       desc, F_BODY, Pt(11), False, C_WHITE)

# 代码统计
tb(s, Cm(21.5), Cm(2.5), Cm(7), Cm(0.6),
   "代码规模", F_BODY, Pt(15), True, C_PRI)
tbl(s, Cm(21.5), Cm(3.3), Cm(7), Cm(4),
    [["模块", "行数"],
     ["ROS2 C++", "6,990"],
     ["配置 YAML", "840"],
     ["CMake+Launch", "1,248"],
     ["STM32 嵌入式", "~175K"],
     ["合计", "~184K"]],
    [Cm(4), Cm(3)], Pt(11))

mpara(s, Cm(21.5), Cm(7.8), Cm(7), Cm(1.5), [
    {"text": "✅ 上层可依赖下层", "fs": Pt(11), "color": C_GREEN},
    {"text": "❌ 禁止反向依赖", "fs": Pt(11), "color": C_RED},
    {"text": "22个ROS2依赖包 · 13个编译目标", "fs": Pt(10), "color": C_GRAY},
])


# ==========================================================
# P7 - 嵌入式-运动学算法
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "嵌入式 — 运动学算法")
section_tag(s, "Part 3 · 进展汇报")
page_num(s, 7)

placeholder(s, ML, Cm(2.5), Cm(12), Cm(5.5),
            "[D-H坐标系图 🎨2\n标注6个关节轴和连杆参数]")

# 位置IK
rrect(s, Cm(15.5), Cm(2.5), Cm(7), Cm(3.8), RGBColor(0xE3, 0xF2, 0xFD), border=RGBColor(0x1E, 0x88, 0xE5))
mpara(s, Cm(15.8), Cm(2.6), Cm(6.4), Cm(3.5), [
    {"text": "位置IK (J1-J3)", "fs": Pt(14), "bold": True, "color": RGBColor(0x1E, 0x88, 0xE5)},
    {"text": "方法: 解析几何法", "fs": Pt(12)},
    {"text": "1. 求腕部中心 Pw=Pee-d6·R·ẑ", "fs": Pt(11)},
    {"text": "2. 几何投影求 θ1,θ2,θ3", "fs": Pt(11)},
    {"text": "每组≤2个解", "fs": Pt(10), "color": C_GRAY},
])

# 姿态IK
rrect(s, Cm(23), Cm(2.5), Cm(7), Cm(3.8), RGBColor(0xE8, 0xF5, 0xE9), border=RGBColor(0x2E, 0x7D, 0x32))
mpara(s, Cm(23.3), Cm(2.6), Cm(6.4), Cm(3.5), [
    {"text": "姿态IK (J4-J6)", "fs": Pt(14), "bold": True, "color": RGBColor(0x2E, 0x7D, 0x32)},
    {"text": "方法: 旋转矩阵分解", "fs": Pt(12)},
    {"text": "1. R₃₆ = (R₀₃)⁻¹ · R₀₆", "fs": Pt(11)},
    {"text": "2. ZYZ欧拉角提取θ4,θ5,θ6", "fs": Pt(11)},
    {"text": "atan2求解", "fs": Pt(10), "color": C_GRAY},
])

# 底部性能
rrect(s, Cm(15.5), Cm(6.8), Cm(14.5), Cm(1), C_PRI,
      "总计算时间 < 50μs @ STM32F4 · 满足1000Hz实时需求",
      Pt(13), C_WHITE, True, border=None)

# 公式
tb(s, ML, Cm(8.5), Cm(12), Cm(1.5),
   "Pw = Pee - d₆ · Rtarget · ẑ        θ₁ = atan2(Pwy, Pwx)",
   F_MATH, Pt(14), False, C_DARK)


# ==========================================================
# P8 - 嵌入式-状态机 (mermaid FSM图)
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "嵌入式 — 状态机与控制")
section_tag(s, "Part 3 · 进展汇报")
page_num(s, 8)

add_mermaid(s, "fsm", Cm(1), Cm(2.2), Cm(17), Cm(9))

bullets(s, Cm(19), Cm(2.5), Cm(10), Cm(8), [
    "1ms主循环 (1000Hz硬实时)",
    "",
    "C++模板化接口，易扩展新状态",
    "",
    "多协议电机统一驱动:",
    "  · DJI (C620/C610)",
    "  · 达妙 (DM-J8009)",
    "  · 瓴控 (LK Motor)",
    "",
    "任意状态可触发紧急复位(Escape)",
], Pt(13), C_BLACK, Pt(2))


# ==========================================================
# P9 - 嵌入式-通信与集成
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "嵌入式 — 通信与集成")
section_tag(s, "Part 3 · 进展汇报")
page_num(s, 9)

tb(s, ML, Cm(2.3), Cm(15), Cm(0.6),
   "Seasky串口协议 (USB-UART, 921600 baud)", F_BODY, Pt(15), True, C_PRI)

# 帧结构表
tbl(s, ML, Cm(3.2), Cm(28), Cm(1.6),
    [["SOF", "DataLen", "CRC8", "CmdID", "Flags", "Payload", "CRC16"],
     ["0xA5", "2B", "1B", "2B", "2B", "NB", "2B"]],
    [Cm(3.5), Cm(4), Cm(3.5), Cm(4), Cm(3.5), Cm(5.5), Cm(4)], Pt(13))

tb(s, ML, Cm(5), Cm(20), Cm(0.4),
   "│← Header(4B) →│← Data →│← CRC →│   CRC8校验头 + CRC16全帧校验",
   F_CODE, Pt(10), False, C_GRAY)

# 数据包类型
tbl(s, ML, Cm(5.8), Cm(28), Cm(5),
    [["CmdID", "方向", "频率", "内容"],
     ["0x0001", "MCU→PC", "200Hz", "关节反馈 (7×pos+vel+alive) 84B"],
     ["0x0002", "PC→MCU", "200Hz", "6轴力矩指令 24B"],
     ["0x0004", "PC→MCU", "50Hz", "夹爪动作 1B (GRIP/RELEASE/STOP)"],
     ["0x0005", "MCU→PC", "按需", "任务指令 3B (cmd+param+seq)"],
     ["0x0006", "PC→MCU", "10Hz", "状态反馈 4B (state+progress+error+gripper)"]],
    [Cm(3.5), Cm(4), Cm(3), Cm(17.5)], Pt(12))


# ==========================================================
# P10 - ROS2节点拓扑 ★ (mermaid图)
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "ROS2 节点拓扑 ★")
section_tag(s, "Part 3 · 进展汇报")
page_num(s, 10)

add_mermaid(s, "ros2_topology", Cm(1.5), Cm(2.2), Cm(20), Cm(9))

mpara(s, Cm(22.5), Cm(2.5), Cm(7), Cm(8), [
    {"text": "7个节点 · 200Hz控制环", "fs": Pt(13), "bold": True, "color": C_PRI},
    {"text": ""},
    {"text": "话题:", "fs": Pt(12), "bold": True},
    {"text": "/joint_states (7关节)", "fs": Pt(10), "font": F_CODE},
    {"text": "/effort_commands (7元素)", "fs": Pt(10), "font": F_CODE},
    {"text": ""},
    {"text": "服务:", "fs": Pt(12), "bold": True},
    {"text": "/gripper_control", "fs": Pt(10), "font": F_CODE},
    {"text": "/list_trajectories", "fs": Pt(10), "font": F_CODE},
    {"text": "/move_to_cartesian_rpy", "fs": Pt(10), "font": F_CODE},
    {"text": ""},
    {"text": "Action:", "fs": Pt(12), "bold": True},
    {"text": "/follow_joint_trajectory", "fs": Pt(10), "font": F_CODE},
])


# ==========================================================
# P11 - 控制算法-动力学前馈 ★★ (合并Kalman)
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "控制算法 — 动力学前馈 ★★")
section_tag(s, "Part 3 · 进展汇报")
page_num(s, 11)

# 核心公式
rrect(s, ML, Cm(2.3), CW, Cm(1.8), RGBColor(0xFC, 0xF5, 0xE0), border=None)
mpara(s, ML + Cm(0.5), Cm(2.4), CW - Cm(1), Cm(1.6), [
    {"text": "τ = M(q)q̈ + C(q,q̇) + G(q) + τ_PID",
     "font": F_MATH, "fs": Pt(22), "bold": True, "color": C_DARK, "align": PP_ALIGN.CENTER},
    {"text": "    惯性力      科氏+离心     重力      PID反馈",
     "fs": Pt(11), "color": C_GRAY, "align": PP_ALIGN.CENTER},
])

# 左侧: 实现 + 4模式
mpara(s, ML, Cm(4.5), Cm(14), Cm(3.5), [
    {"text": "实现", "fs": Pt(14), "bold": True, "color": C_PRI},
    {"text": "· KDL::ChainDynParam 求解M/C/G矩阵", "fs": Pt(12)},
    {"text": "· URDF自动解析关节链 · 计算 <1ms @200Hz", "fs": Pt(12)},
    {"text": "· 编码器速度信号经后处理滤波平滑 (消除高频噪声)", "fs": Pt(12), "color": RGBColor(0xE6, 0x5C, 0x00)},
])

tbl(s, ML, Cm(7.5), Cm(14), Cm(3.5),
    [["模式", "说明"],
     ["RELAX", "零力矩, 电机断力 (上电默认)"],
     ["FREEDRIVE", "仅重力补偿G(q), 可拖动教学"],
     ["OVERDRIVE", "G(q)+PD反馈, 锁位保持"],
     ["EXECUTE", "全前馈+PID, 轨迹执行中"]],
    [Cm(4), Cm(10)], Pt(12))

# 右侧: 控制框图说明
mpara(s, Cm(17), Cm(4.5), Cm(12), Cm(6.5), [
    {"text": "控制框图", "fs": Pt(14), "bold": True, "color": C_PRI},
    {"text": ""},
    {"text": "q_target → MoveIt2轨迹生成", "fs": Pt(12), "font": F_CODE},
    {"text": ""},
    {"text": "τ_ff = M(q)q̈ + C(q,q̇) + G(q)", "fs": Pt(13), "font": F_MATH, "bold": True},
    {"text": "    → 前馈力矩 (消除重力/惯性)", "fs": Pt(11), "color": C_GRAY},
    {"text": ""},
    {"text": "τ_fb = 级联PID(e_pos, e_vel)", "fs": Pt(13), "font": F_MATH, "bold": True},
    {"text": "    → 反馈校正 (消除跟踪误差)", "fs": Pt(11), "color": C_GRAY},
    {"text": ""},
    {"text": "τ_total = τ_ff + τ_fb → 电机", "fs": Pt(14), "bold": True, "color": C_PRI},
])


# ==========================================================
# P12 - 级联PID ★★ (mermaid框图)
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "控制算法 — 级联PID ★★")
section_tag(s, "Part 3 · 进展汇报")
page_num(s, 12)

# mermaid PID框图
add_mermaid(s, "cascade_pid", Cm(1), Cm(2.3), Cm(28), Cm(3.5))

# 公式
mpara(s, ML, Cm(6), Cm(28), Cm(1.2), [
    {"text": "外环: v_ref = clamp(Kp·e_pos + Ki·∫e_pos, -v_max, +v_max)        内环: τ_fb = Kp_vel·e_vel + Ki_vel·∫e_vel",
     "font": F_MATH, "fs": Pt(12), "color": C_DARK},
])

# 参数表
tbl(s, ML, Cm(7.5), Cm(17), Cm(4),
    [["关节", "pos_Kp", "pos_Ki", "vel_Kp", "vel_Ki", "vel_lim"],
     ["J1", "4.0", "0.3", "7.0", "0.5", "10 rad/s"],
     ["J2", "6.0", "0.3", "7.0", "0.5", "10 rad/s"],
     ["J3", "4.0", "0.3", "6.0", "0.5", "10 rad/s"],
     ["J4", "1.5", "0.2", "6.3", "0.3", "10 rad/s"],
     ["J5", "1.5", "0.2", "1.0", "0.1", "10 rad/s"],
     ["J6", "1.0", "0.15", "1.0", "0.1", "10 rad/s"]],
    [Cm(2.5), Cm(2.5), Cm(2.5), Cm(2.5), Cm(2.5), Cm(4.5)], Pt(11))

# 设计要点
bullets(s, Cm(19.5), Cm(7.5), Cm(10), Cm(4), [
    "近端(J1-3): 高增益 ← J8009大力矩",
    "远端(J4-6): 低增益 ← 小电机",
    "三重抗积分饱和:",
    "  条件积分(e<0.1rad)",
    "  积分限幅(max_integral)",
    "  急停清零",
    "所有Kd=0: 纯P+PI结构",
], Pt(11), C_BLACK, Pt(2), "·")


# ==========================================================
# P13 - MuJoCo仿真与数字孪生 (原P14，Kalman页被删)
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "MuJoCo仿真与数字孪生")
section_tag(s, "Part 3 · 进展汇报")
page_num(s, 13)

placeholder(s, ML, Cm(2.5), Cm(15), Cm(8.5),
            "[MuJoCo仿真窗口截图 🖥1\n含ImGui参数面板 + ImPlot实时曲线叠加]")

mpara(s, Cm(17.5), Cm(2.5), Cm(12), Cm(4.5), [
    {"text": "① 物理仿真模式", "fs": Pt(15), "bold": True, "color": RGBColor(0x1E, 0x88, 0xE5)},
    {"text": "· MuJoCo物理引擎", "fs": Pt(12)},
    {"text": "· 力矩输入 → 物理响应 · 碰撞检测", "fs": Pt(12)},
    {"text": "· 算法开发调试", "fs": Pt(12)},
    {"text": ""},
    {"text": "② 数字孪生模式", "fs": Pt(15), "bold": True, "color": RGBColor(0x2E, 0x7D, 0x32)},
    {"text": "· visualization_only=true", "fs": Pt(12), "font": F_CODE},
    {"text": "· 真机数据驱动 · 实时3D同步", "fs": Pt(12)},
    {"text": "· 现场监控调试", "fs": Pt(12)},
])

mpara(s, Cm(17.5), Cm(7.5), Cm(12), Cm(3.5), [
    {"text": "可视化特性", "fs": Pt(14), "bold": True, "color": C_PRI},
    {"text": "· ImGui参数面板 (实时调参)", "fs": Pt(12)},
    {"text": "· ImPlot实时曲线 (位置/速度/力矩)", "fs": Pt(12)},
    {"text": "· 障碍物场景渲染 (矿石框+能量块)", "fs": Pt(12)},
])


# ==========================================================
# P14 - 应用层
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "应用层 — 轨迹管理与任务执行")
section_tag(s, "Part 3 · 进展汇报")
page_num(s, 14)

card3 = [
    ("Trajectory Manager", "轨迹管理 (518行)", [
        "YAML持久化 保存/加载/执行",
        "ROS2 Action 异步执行",
        "4个Service:",
        " /list · /load · /save · /save_last",
    ], RGBColor(0x42, 0xA5, 0xF5)),
    ("Mission Executor", "任务执行 (1194行)", [
        "ncurses TUI 终端交互",
        "19态状态机 pick_and_deliver",
        "快捷键 [E]执行 [X]复位 [T]轨迹",
        "操控手指令 via CmdID 0x0005",
    ], RGBColor(0x2E, 0x7D, 0x32)),
    ("Cartesian Controller", "笛卡尔控制 (491行)", [
        "Pilz规划器 LIN/PTP · 自动降级",
        "球形工作空间 r∈[0.01,0.8]m",
        "视觉伺服接口预留 20Hz",
        "Service: /move_to_cartesian_rpy",
    ], RGBColor(0xAB, 0x47, 0xBC)),
]
c3w = Cm(9.5)
for i, (t_, sub, items, color) in enumerate(card3):
    x = ML + i * (c3w + Cm(0.3))
    rrect(s, x, Cm(2.5), c3w, Cm(7.5), C_LIGHT, border=color, bw=Pt(2))
    tb(s, x + Cm(0.4), Cm(2.7), c3w - Cm(0.8), Cm(0.6),
       t_, F_CODE, Pt(13), True, color)
    tb(s, x + Cm(0.4), Cm(3.4), c3w - Cm(0.8), Cm(0.5),
       sub, F_BODY, Pt(11), False, C_GRAY)
    bullets(s, x + Cm(0.4), Cm(4.2), c3w - Cm(0.8), Cm(5),
            items, Pt(11), C_BLACK, Pt(3), "·")


# ==========================================================
# P15 - Sim2Real全链路 ★ (mermaid流程)
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "Sim2Real — 仿真→真机全链路 ★")
section_tag(s, "Part 3 · 进展汇报")
page_num(s, 15)

add_mermaid(s, "sim2real", Cm(3), Cm(2.5), Cm(24), Cm(2.5))

# 照片/截图占位
placeholder(s, ML, Cm(5.5), Cm(14), Cm(4),
            "[真机运行照 📷7\n机器人执行矿石搬运动作]")
placeholder(s, Cm(16.5), Cm(5.5), Cm(13), Cm(4),
            "[数字孪生截图 🖥2\nMuJoCo实时同步真机姿态]")

# 代码规模
tbl(s, ML, Cm(10), Cm(28), Cm(1.5),
    [["ROS2 C++ 源码", "配置 YAML", "CMake+Launch", "STM32 嵌入式", "合计"],
     ["6,990行 (16文件)", "840行", "1,248行", "~175K行", "~184K行"]],
    [Cm(6), Cm(5), Cm(5), Cm(6), Cm(6)], Pt(11))


# ==========================================================
# P16 - 实验数据 ★
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "实验验证数据 ★")
section_tag(s, "Part 3 · 进展汇报")
page_num(s, 16)

placeholder(s, ML, Cm(2.5), Cm(14), Cm(4.5),
            "[关节跟踪精度 🖥3-a\n蓝:目标  红:实际  绿:误差]")
placeholder(s, Cm(16.5), Cm(2.5), Cm(13), Cm(4.5),
            "[力矩响应对比 🖥3-b\n蓝:有前馈  红:纯PID]")

tbl(s, ML, Cm(7.5), Cm(13.5), Cm(3.5),
    [["指标", "实测值"],
     ["控制频率", "200Hz (5ms周期)"],
     ["控制回路延迟", "< 5ms"],
     ["动力学计算耗时", "< 1ms"],
     ["IK计算(STM32)", "< 50μs"],
     ["串口波特率", "921,600 baud"]],
    [Cm(6), Cm(7.5)], Pt(12))

tbl(s, Cm(16.5), Cm(7.5), Cm(13), Cm(3.5),
    [["安全保护机制", "触发条件"],
     ["传感器超时", "100ms无数据→急停"],
     ["位置误差急停", "> 0.8 rad触发"],
     ["速度越限拒绝", "> 20 rad/s丢弃"],
     ["力矩限幅", "各关节独立"],
     ["NaN/Inf检测", "实时检查→零力矩"]],
    [Cm(5.5), Cm(7.5)], Pt(12))


# ==========================================================
# P17 - 遇到的问题与解决
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "遇到的问题与解决")
section_tag(s, "Part 4 · 问题与计划")
page_num(s, 17)

problems = [
    ("① 启动关节漂移",
     "电机上电瞬间力矩为零，\n关节在重力下下坠",
     "首帧位置自动锁定为目标\n(HOLD模式)，立即施加G(q)"),
    ("② 力矩输出震荡",
     "编码器速度信号噪声大，\nPID输出严重抖动",
     "速度信号后处理滤波平滑\n(参数在线可调)"),
    ("③ MuJoCo黑屏",
     "多线程同时访问OpenGL\n上下文导致渲染崩溃",
     "OpenGL上下文隔离到\n专用渲染线程"),
]
for i, (t_, prob, sol) in enumerate(problems):
    y = Cm(2.5) + i * Cm(3)
    tb(s, ML, y, Cm(6), Cm(0.6), t_, F_BODY, Pt(15), True, C_RED)

    rrect(s, ML, y + Cm(0.7), Cm(13.5), Cm(2), RGBColor(0xFD, 0xED, 0xED),
           border=RGBColor(0xF5, 0xC6, 0xC6))
    tb(s, ML + Cm(0.3), y + Cm(0.5), Cm(2), Cm(0.5),
       "问题 ❌", F_BODY, Pt(10), True, C_RED)
    tb(s, ML + Cm(0.3), y + Cm(1.1), Cm(12.5), Cm(1.3),
       prob, F_BODY, Pt(12))

    rrect(s, Cm(16.5), y + Cm(0.7), Cm(13), Cm(2), RGBColor(0xE8, 0xF5, 0xE9),
           border=RGBColor(0xC8, 0xE6, 0xC9))
    tb(s, Cm(16.8), y + Cm(0.5), Cm(2), Cm(0.5),
       "解决 ✅", F_BODY, Pt(10), True, C_GREEN)
    tb(s, Cm(16.8), y + Cm(1.1), Cm(12), Cm(1.3),
       sol, F_BODY, Pt(12))


# ==========================================================
# P18 - 后续计划 (mermaid甘特图)
# ==========================================================
s = blank(prs)
add_bg(s)
title(s, "后续工作计划")
section_tag(s, "Part 4 · 问题与计划")
page_num(s, 18)

add_mermaid(s, "gantt_plan", Cm(1.5), Cm(2.5), Cm(27), Cm(5))

# 核心目标
rrect(s, ML, Cm(8), CW, Cm(1.5), C_PRI,
      "核心目标: 打通 视觉检测→位姿估计→运动规划→力矩执行 完整闭环\n(cartesian_controller已预留20Hz接口)",
      Pt(14), C_WHITE, True, border=None)

mpara(s, ML, Cm(10), CW, Cm(1.5), [
    {"text": "3月: 视觉伺服 (YOLOv8+PnP+标定)    4月: 闭环抓取联调    4月底: 论文定稿    5月底: 终期答辩",
     "fs": Pt(13), "color": C_DARK},
])


# ==========================================================
# P19 - 致谢
# ==========================================================
s = blank(prs)
add_bg(s)
page_num(s, 19)

tb(s, Cm(3), Cm(2.5), CW - Cm(3), Cm(2),
   "谢谢各位老师！", F_TITLE, Pt(40), True, C_PRI, PP_ALIGN.CENTER)
tb(s, Cm(3), Cm(5), CW - Cm(3), Cm(1),
   "感谢导师 柯文德 教授的悉心指导", F_BODY, Pt(18), False, C_DARK, PP_ALIGN.CENTER)
tb(s, Cm(3), Cm(6.5), CW - Cm(3), Cm(1),
   "感谢实验室团队的大力支持", F_BODY, Pt(18), False, C_DARK, PP_ALIGN.CENTER)
tb(s, Cm(3), Cm(8.5), CW - Cm(3), Cm(1),
   "欢迎提问与指导", F_BODY, Pt(20), True, C_PRI, PP_ALIGN.CENTER)


# ============================================================
# 保存
# ============================================================
prs.save(OUTPUT_PATH)
print(f"\n✅ PPT已生成: {OUTPUT_PATH}")
print(f"总页数: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.paragraphs[0].text.strip()
            if t and len(t) < 30 and "单击" not in t:
                texts.append(t)
                break
    print(f"  P{i+1}: {texts[0] if texts else '...'}")
